"""
Suplex 14슬라이드 소개서 → PPTX 변환기
- SSOT: docs/sales/JTBD캔버스_P1.md
- 입력 본문: 원페이저.md 기반 + 14슬라이드 구조
- 출력: docs/sales/소개서.pptx
- 16:9, 한국어 (맑은 고딕)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── 디자인 토큰 ──
NAVY_DARK   = RGBColor(0x1e, 0x3a, 0x5f)
NAVY_MID    = RGBColor(0x2d, 0x4a, 0x6f)
WHITE       = RGBColor(0xff, 0xff, 0xff)
GRAY_50     = RGBColor(0xf9, 0xfa, 0xfb)
GRAY_200    = RGBColor(0xe5, 0xe7, 0xeb)
GRAY_400    = RGBColor(0x9c, 0xa3, 0xaf)
GRAY_600    = RGBColor(0x4b, 0x55, 0x63)
GRAY_800    = RGBColor(0x1f, 0x29, 0x37)
RED_700     = RGBColor(0xb9, 0x1c, 0x1c)
GREEN_700   = RGBColor(0x15, 0x80, 0x3d)
AMBER       = RGBColor(0xfb, 0xbf, 0x24)
AMBER_DARK  = RGBColor(0xb4, 0x80, 0x0a)

FONT_KO = "맑은 고딕"
FONT_MONO = "맑은 고딕"
TOTAL = 14


# ── 헬퍼 ──
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(prs, slide, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    spTree = sh._element.getparent()
    spTree.remove(sh._element)
    spTree.insert(2, sh._element)


def text(slide, x, y, w, h, txt, size=14, bold=False, color=GRAY_800,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.name = font or FONT_KO
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def rrect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def footer(slide, num, dark=False):
    color = GRAY_400
    text(slide, Inches(0.5), Inches(7.1), Inches(2), Inches(0.3),
         "SUPLEX", size=9, bold=True, color=color)
    text(slide, Inches(11), Inches(7.1), Inches(1.8), Inches(0.3),
         f"{num:02d} / {TOTAL:02d}", size=9, color=color, align=PP_ALIGN.RIGHT)


def code_block(slide, x, y, w, h, code_text):
    rrect(slide, x, y, w, h, NAVY_DARK)
    tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.2),
                                   w - Inches(0.5), h - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(code_text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = FONT_MONO
        r.font.size = Pt(10)
        r.font.color.rgb = WHITE


def circle_num(slide, x, y, d, num, bg, fg):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.fill.solid(); c.fill.fore_color.rgb = bg
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.name = FONT_KO; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = fg


# ── 슬라이드 ──

def s_1_title(prs):
    s = add_blank(prs); fill_bg(prs, s, NAVY_DARK)
    rect(s, Inches(0.6), Inches(2.2), Inches(0.08), Inches(3.4), AMBER)
    text(s, Inches(0.95), Inches(2.2), Inches(11), Inches(0.8),
         "수플렉스", size=64, bold=True, color=WHITE)
    text(s, Inches(0.95), Inches(3.25), Inches(11), Inches(0.4),
         "Suplex", size=18, color=GRAY_400)
    text(s, Inches(0.95), Inches(4.0), Inches(12), Inches(0.7),
         "한 사람의 능력이 아니라,", size=26, bold=True, color=GRAY_200)
    text(s, Inches(0.95), Inches(4.65), Inches(12), Inches(0.7),
         "회사 전체의 능력이 올라갑니다.", size=26, bold=True, color=AMBER)
    text(s, Inches(0.95), Inches(5.6), Inches(11), Inches(0.5),
         "효율로 더 많은 매출을 만드는 인테리어 회사 운영 도구",
         size=14, color=GRAY_200)
    text(s, Inches(0.5), Inches(7.05), Inches(12), Inches(0.4),
         "도입 검토용 소개서  ·  2026", size=10, color=GRAY_400, align=PP_ALIGN.RIGHT)


def s_2_pain(prs):
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "01. PAIN", size=11, bold=True, color=GRAY_400)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "이런 적 있으세요?", size=30, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.85), Inches(0.6), Inches(0.04), AMBER)

    pains = [
        ('현장에서 사무실에 전화',
         '"마감재 뭐였지?"  "타일 언제 잡았지?"\n"마감재 확정됐어?"'),
        ('작업자·자재상 섭외 십수번 반복',
         '현장 주소·상황을 매번 같은 내용으로\n손으로 타이핑'),
        ('프로젝트 노하우 휘발',
         '끝난 프로젝트의 피드백·자재 노하우 —\n어디 남겼는지 까먹어서 반영 X'),
        ('발주 데드라인 망실',
         '"오늘까지 발주했어야 했는데" —\n사람이 까먹어서 자재가 늦음'),
    ]
    cw = Inches(6.0); ch = Inches(1.95); gx = Inches(0.15); gy = Inches(0.15)
    sx = Inches(0.7); sy = Inches(2.2)
    for i, (label, body) in enumerate(pains):
        col, row = i % 2, i // 2
        x = sx + (cw + gx) * col
        y = sy + (ch + gy) * row
        rrect(s, x, y, cw, ch, GRAY_50)
        text(s, x + Inches(0.35), y + Inches(0.3), cw - Inches(0.7), Inches(0.45),
             label, size=15, bold=True, color=NAVY_DARK)
        text(s, x + Inches(0.35), y + Inches(0.85), cw - Inches(0.7), Inches(1.0),
             body, size=12, color=GRAY_800)

    rrect(s, Inches(0.7), Inches(6.4), Inches(12.15), Inches(0.55), NAVY_DARK)
    text(s, Inches(0.7), Inches(6.5), Inches(12.15), Inches(0.4),
         "기록이 한곳에 정렬되지 않으면 놓치는 정보가 생기기 마련입니다.",
         size=12, bold=True, color=AMBER,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    footer(s, 2)


def s_3_old_vs_new(prs):
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "02. CONTRAST", size=11, bold=True, color=GRAY_400)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "예전 방식 vs 수플렉스", size=30, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.85), Inches(0.6), Inches(0.04), AMBER)

    # 좌측: 예전 방식
    rrect(s, Inches(0.7), Inches(2.2), Inches(6.0), Inches(4.6), GRAY_50)
    text(s, Inches(1.0), Inches(2.4), Inches(5.5), Inches(0.4),
         "예전 방식", size=12, bold=True, color=GRAY_400)
    text(s, Inches(1.0), Inches(2.85), Inches(5.5), Inches(0.6),
         "엑셀 + 카톡 + 종이노트", size=20, bold=True, color=GRAY_800)
    old_items = [
        "정보가 3~5곳에 흩어짐",
        "변경 시 매번 수동 동기화",
        "현장에서 사무실에 전화로 확인",
        "작업자 안내 매번 손으로 타이핑",
        "데드라인은 사람이 기억",
        "노하우는 한 사람의 머릿속에",
    ]
    y = Inches(3.75)
    for item in old_items:
        text(s, Inches(1.2), y, Inches(0.3), Inches(0.4),
             "✕", size=14, bold=True, color=RED_700)
        text(s, Inches(1.55), y, Inches(5.0), Inches(0.4),
             item, size=12, color=GRAY_800)
        y += Inches(0.45)

    # 우측: 수플렉스
    rrect(s, Inches(6.85), Inches(2.2), Inches(6.0), Inches(4.6), NAVY_DARK)
    text(s, Inches(7.15), Inches(2.4), Inches(5.5), Inches(0.4),
         "수플렉스", size=12, bold=True, color=AMBER)
    text(s, Inches(7.15), Inches(2.85), Inches(5.5), Inches(0.6),
         "한 화면, 시스템이 챙김", size=20, bold=True, color=WHITE)
    new_items = [
        "한 화면에 통합 (공정 현황 통합 뷰)",
        "한 곳 입력 → 자동 연결",
        "모바일에서 즉시 확인",
        "카톡 텍스트 시스템이 자동 정리",
        "시스템이 D-day 자동 챙김",
        "노하우는 회사 자산으로 누적",
    ]
    y = Inches(3.75)
    for item in new_items:
        text(s, Inches(7.35), y, Inches(0.3), Inches(0.4),
             "✓", size=14, bold=True, color=AMBER)
        text(s, Inches(7.7), y, Inches(5.0), Inches(0.4),
             item, size=12, color=WHITE)
        y += Inches(0.45)

    footer(s, 3)


def s_4_solution_overview(prs):
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "03. SOLUTION", size=11, bold=True, color=GRAY_400)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "수플렉스가 잡습니다", size=30, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.85), Inches(0.6), Inches(0.04), AMBER)

    sols = [
        ("01", "지난 프로젝트의 노하우가 지금 살아납니다",
         "프로젝트 기록 검색 + 가전 DB 누적 + 자재 불러오기", False),
        ("02", "시스템이 데드라인을 챙깁니다",
         "한국 표준 5묶음 자동 적용 + D-day 자동 체크리스트", False),
        ("03", "카톡 텍스트 자동 정리",
         "일정 복사 + 발주서 자동 작성 + 인건비 정산 카톡", True),
        ("04", "4~10개 현장이 한 화면에서",
         "공정 현황 통합 뷰 + 모바일 가독성", True),
    ]
    cw = Inches(6.0); ch = Inches(2.1); gx = Inches(0.15); gy = Inches(0.15)
    sx = Inches(0.7); sy = Inches(2.3)
    for i, (num, title, desc, hl) in enumerate(sols):
        col, row = i % 2, i // 2
        x = sx + (cw + gx) * col
        y = sy + (ch + gy) * row
        bg = NAVY_DARK if hl else GRAY_50
        c_n = AMBER if hl else AMBER_DARK
        c_t = WHITE if hl else NAVY_DARK
        c_d = GRAY_200 if hl else GRAY_600
        rrect(s, x, y, cw, ch, bg)
        text(s, x + Inches(0.4), y + Inches(0.25), cw - Inches(0.8), Inches(0.5),
             num, size=28, bold=True, color=c_n)
        text(s, x + Inches(0.4), y + Inches(0.9), cw - Inches(0.8), Inches(0.5),
             title, size=15, bold=True, color=c_t)
        text(s, x + Inches(0.4), y + Inches(1.45), cw - Inches(0.8), Inches(0.55),
             desc, size=11, color=c_d)
        if hl:
            text(s, x + cw - Inches(1.2), y + Inches(0.3), Inches(1.0), Inches(0.3),
                 "MAIN", size=9, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)
    footer(s, 4)


def s_5_solution_1(prs):
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "03. SOLUTION  ·  ①", size=11, bold=True, color=GRAY_400)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "지난 프로젝트의 노하우가 지금 현장에서 다시 살아납니다",
         size=22, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.95), Inches(0.6), Inches(0.04), AMBER)

    rrect(s, Inches(0.7), Inches(2.3), Inches(7.8), Inches(1.3), GRAY_50, line=AMBER)
    text(s, Inches(1.0), Inches(2.55), Inches(7.5), Inches(0.4),
         "시나리오", size=11, bold=True, color=AMBER_DARK)
    text(s, Inches(1.0), Inches(2.95), Inches(7.5), Inches(0.5),
         '"3년 전 그 화이트 강마루, 입주 후 반응이 어땠지?"',
         size=15, bold=True, color=GRAY_800)

    body_lines = [
        "검색 한 번으로 그때의 피드백·A/S 기록·고객 반응이 떠오릅니다.",
        "비슷한 사례의 현장이 들어왔을 때, 회사가 쌓아온 결정과",
        "결과가 그대로 다음 판단에 쓰입니다.",
        "",
        "견적·자재·가전 모델은 한 번 입력하면 회사 데이터에 누적,",
        "두 번째 같은 자재를 처음부터 만들 필요가 없습니다.",
    ]
    y = Inches(3.85)
    for line in body_lines:
        text(s, Inches(0.7), y, Inches(8.0), Inches(0.4),
             line, size=13, color=GRAY_800)
        y += Inches(0.4)

    code = ("[검색: 강마루]\n\n2024-03  강남\n동부센트레빌\n303-1502\n\n"
            "자재 : 동화 화이트\n       강마루\n반응 : 우수\n       2년차 변색 X\n"
            "A/S  : 1회\n       걸레받이 코킹\n\n→ 같은 자재\n  추천 가능")
    code_block(s, Inches(8.85), Inches(2.3), Inches(4.0), Inches(4.4), code)
    footer(s, 5)


def s_6_solution_2(prs):
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "03. SOLUTION  ·  ②", size=11, bold=True, color=GRAY_400)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "시스템이 데드라인을 챙깁니다",
         size=26, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.95), Inches(0.6), Inches(0.04), AMBER)

    rrect(s, Inches(0.7), Inches(2.3), Inches(7.8), Inches(1.55), NAVY_DARK)
    text(s, Inches(1.0), Inches(2.5), Inches(7.5), Inches(0.4),
         "가입 첫날부터", size=11, bold=True, color=AMBER)
    text(s, Inches(1.0), Inches(2.85), Inches(7.5), Inches(0.5),
         "한국 인테리어 표준 5묶음이", size=18, bold=True, color=WHITE)
    text(s, Inches(1.0), Inches(3.3), Inches(7.5), Inches(0.5),
         "회사에 깔려있습니다", size=18, bold=True, color=WHITE)

    text(s, Inches(0.7), Inches(4.05), Inches(8), Inches(0.4),
         "공정 · 키워드 · D-N 룰 · 체크리스트 · 견적 가이드",
         size=11, color=GRAY_600)

    text(s, Inches(0.7), Inches(4.7), Inches(8), Inches(0.4),
         '일정에 "도배" 입력 →', size=14, bold=True, color=NAVY_DARK)
    rows = [
        ("D-3", '"벽지 도착 확인" 자동 추가'),
        ("D-1", '"초벌 풀칠 준비" 자동 추가'),
    ]
    y = Inches(5.2)
    for label, desc in rows:
        rrect(s, Inches(0.95), y, Inches(0.7), Inches(0.45), AMBER)
        text(s, Inches(0.95), y + Inches(0.05), Inches(0.7), Inches(0.4),
             label, size=12, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
        text(s, Inches(1.85), y + Inches(0.05), Inches(7), Inches(0.45),
             desc, size=14, color=GRAY_800)
        y += Inches(0.55)

    rrect(s, Inches(8.85), Inches(2.3), Inches(4.0), Inches(4.4), AMBER)
    text(s, Inches(8.95), Inches(3.4), Inches(3.8), Inches(0.5),
         "사람이 까먹어도", size=15, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
    text(s, Inches(8.95), Inches(3.95), Inches(3.8), Inches(0.5),
         "시스템이", size=15, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
    text(s, Inches(8.95), Inches(4.5), Inches(3.8), Inches(0.5),
         "까먹지 않습니다", size=15, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
    footer(s, 6)


def s_7_solution_3(prs):
    """메인 어필 컷"""
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "03. SOLUTION  ·  ③  ·  메인 기능", size=11, bold=True, color=AMBER_DARK)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "카톡으로 보낼 텍스트는 시스템이 정리합니다",
         size=22, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.95), Inches(0.6), Inches(0.04), AMBER)

    items = [
        ("일정 복사",
         '키워드 "전기" 한 번 →\n4개 현장 일정·주소·특이사항이\n카톡 텍스트로 한 번에 복사'),
        ("발주서 자동 작성",
         "마감재 체크 →\n자재상에게 보낼 카톡 텍스트로\n자동 변환"),
        ("인건비 정산",
         "작업자별 [일수×단가\n+식비+교통비] 자동 합산 →\n송금용 카톡 텍스트"),
    ]
    cw = Inches(4.0); gx = Inches(0.075)
    sx = Inches(0.7); sy = Inches(2.3); ch = Inches(2.55)
    for i, (label, body) in enumerate(items):
        x = sx + (cw + gx) * i
        rrect(s, x, sy, cw, ch, NAVY_DARK)
        rect(s, x + Inches(0.3), sy + Inches(0.3), Inches(0.5), Inches(0.04), AMBER)
        text(s, x + Inches(0.3), sy + Inches(0.45), cw - Inches(0.6), Inches(0.45),
             label, size=15, bold=True, color=AMBER)
        text(s, x + Inches(0.3), sy + Inches(1.05), cw - Inches(0.6), Inches(1.4),
             body, size=12, color=WHITE)

    rrect(s, Inches(0.7), Inches(5.15), Inches(12.15), Inches(1.0), AMBER)
    text(s, Inches(0.7), Inches(5.3), Inches(12.15), Inches(0.5),
         "매번 십수번 반복 타이핑하던 잡일이",
         size=18, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(5.7), Inches(12.15), Inches(0.5),
         "한 번의 클릭으로 끝납니다.",
         size=18, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)

    text(s, Inches(0.7), Inches(6.4), Inches(12.15), Inches(0.4),
         "* 카카오톡 알림 준비 중 — 베타 이후 회사의 알림 채널로 임박한 데드라인이 자동 전달됩니다.",
         size=10, color=GRAY_600, align=PP_ALIGN.CENTER)
    footer(s, 7)


def s_8_solution_4(prs):
    """메인 어필 컷"""
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "03. SOLUTION  ·  ④  ·  메인 기능", size=11, bold=True, color=AMBER_DARK)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "4~10개 현장이 한 화면에서 보입니다",
         size=22, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.95), Inches(0.6), Inches(0.04), AMBER)

    text(s, Inches(0.7), Inches(2.4), Inches(7.5), Inches(0.6),
         "공정 현황 통합 뷰", size=22, bold=True, color=AMBER_DARK)

    body = [
        "25공정 × (견적·마감재·일정·발주) 매트릭스",
        "빠뜨린 항목·임박한 데드라인이 색깔로 강조",
        "프로젝트별로 어디까지 왔는지 한눈에",
        "",
        "모바일에서 글씨가 직접 보입니다",
        "팝업 클릭으로 확인하던 잡일이 사라집니다",
        "현장에서 사무실에 전화할 일이 없어집니다",
    ]
    y = Inches(3.15)
    for line in body:
        if not line:
            y += Inches(0.15); continue
        text(s, Inches(0.95), y, Inches(7.5), Inches(0.4),
             "▸ " + line, size=13, color=GRAY_800)
        y += Inches(0.4)

    # 우측: 미니 매트릭스 시각화
    rrect(s, Inches(8.6), Inches(2.4), Inches(4.25), Inches(4.4), GRAY_50)
    text(s, Inches(8.85), Inches(2.6), Inches(3.8), Inches(0.4),
         "공정 현황 (예시)", size=11, bold=True, color=GRAY_600)

    headers = ["공정", "견", "재", "일", "발"]
    col_w = [Inches(1.4), Inches(0.55), Inches(0.55), Inches(0.55), Inches(0.55)]
    hx = Inches(8.85); hy = Inches(3.15)
    cum = Inches(0)
    for i, h in enumerate(headers):
        text(s, hx + cum, hy, col_w[i], Inches(0.3),
             h, size=10, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
        cum = cum + col_w[i]

    rows = [
        ("타일", "✓", "✓", "✓", "!"),
        ("도배", "✓", "✓", "✗", "✗"),
        ("전기", "✓", "✓", "✓", "✓"),
        ("바닥", "✓", "✓", "!", "✗"),
        ("주방", "✓", "!", "✗", "✗"),
        ("가구", "!", "✗", "✗", "✗"),
    ]
    ry = Inches(3.55)
    for phase, *vals in rows:
        text(s, hx, ry, col_w[0], Inches(0.32),
             phase, size=10, color=GRAY_800)
        cum2 = col_w[0]
        for i, v in enumerate(vals):
            color = (GREEN_700 if v == "✓"
                     else AMBER_DARK if v == "!"
                     else RED_700)
            text(s, hx + cum2, ry, col_w[i+1], Inches(0.32),
                 v, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
            cum2 = cum2 + col_w[i+1]
        ry += Inches(0.42)

    text(s, Inches(8.85), Inches(6.25), Inches(3.8), Inches(0.35),
         "✓ 완료   ! 임박   ✗ 누락",
         size=9, color=GRAY_600, align=PP_ALIGN.CENTER)
    footer(s, 8)


def s_9_company_asset(prs):
    s = add_blank(prs); fill_bg(prs, s, NAVY_DARK)
    text(s, Inches(0.7), Inches(2.4), Inches(12), Inches(0.5),
         "수플렉스의 메인 테마", size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(3.1), Inches(12), Inches(1),
         "회사의 노하우는", size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(4.0), Inches(12), Inches(1),
         "회사 안에 남습니다.", size=42, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.4),
         "견적 기준 · 자재 결정 · 현장 환경 · 고객 대응",
         size=13, color=GRAY_400, align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.4),
         "한 사람의 머릿속이 아니라 회사의 자산으로 누적됩니다.",
         size=13, color=GRAY_200, align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.4),
         "담당자가 바뀌어도 다음 사람이 그 자리를 그대로 이어받습니다.",
         size=13, color=GRAY_200, align=PP_ALIGN.CENTER)
    footer(s, 9, dark=True)


def s_10_diff(prs):
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "04. DIFFERENTIATOR", size=11, bold=True, color=GRAY_400)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "다른 도구와 무엇이 다른가",
         size=30, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.85), Inches(0.6), Inches(0.04), AMBER)

    diffs = [
        ("회사 자산화",
         "직원 머릿속 노하우가 회사 자산으로. 일을 쌓을수록 도구가 더 강해집니다."),
        ("시스템 프리셋",
         "가입 첫날부터 한국 인테리어 표준 5묶음이 깔린 도구. 빈 도구가 아닙니다."),
        ("모바일 우선",
         "현장에서 글씨 직접 보임. 사무실에 전화할 필요가 없습니다."),
        ("클라이언트 비접근",
         "회사 영업 정보가 클라이언트에 노출 X. PDF만 발송됩니다."),
        ("락인 없음",
         "JSON으로 회사 전체 데이터를 통째 백업. 떠나고 싶을 때 떠날 수 있습니다."),
    ]
    # 첫 번째: 풀와이드 강조
    rrect(s, Inches(0.7), Inches(2.2), Inches(12.15), Inches(1.3), NAVY_DARK)
    text(s, Inches(1.1), Inches(2.45), Inches(11.5), Inches(0.45),
         diffs[0][0], size=20, bold=True, color=AMBER)
    text(s, Inches(1.1), Inches(3.0), Inches(11.5), Inches(0.5),
         diffs[0][1], size=14, color=GRAY_200)

    cw = Inches(6.0); ch = Inches(1.45); gx = Inches(0.15); gy = Inches(0.15)
    sx = Inches(0.7); sy = Inches(3.65)
    for i, (title_, desc) in enumerate(diffs[1:]):
        col, row = i % 2, i // 2
        x = sx + (cw + gx) * col
        y = sy + (ch + gy) * row
        rrect(s, x, y, cw, ch, GRAY_50)
        text(s, x + Inches(0.35), y + Inches(0.22), cw - Inches(0.7), Inches(0.4),
             title_, size=15, bold=True, color=NAVY_DARK)
        text(s, x + Inches(0.35), y + Inches(0.7), cw - Inches(0.7), Inches(0.6),
             desc, size=11, color=GRAY_600)
    footer(s, 10)


def s_11_beta(prs):
    s = add_blank(prs); fill_bg(prs, s, NAVY_DARK)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "05. CLOSE BETA", size=11, bold=True, color=AMBER)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "클로즈 베타 — 평생 20% 할인",
         size=30, bold=True, color=WHITE)
    rect(s, Inches(0.7), Inches(1.85), Inches(0.6), Inches(0.04), AMBER)

    items = [
        ("2개월 무료", "결제 정보 등록 X · 카드 등록 X"),
        ("프로 등급 풀 기능", "도구의 모든 기능을 그대로 체험"),
        ("정식 출시 후 평생 20% 할인", "어느 등급이든 영구 적용"),
        ("락인 없음", "JSON 백업으로 데이터 통째 가져갈 수 있음"),
    ]
    cw = Inches(6.0); ch = Inches(1.5); gx = Inches(0.15); gy = Inches(0.15)
    sx = Inches(0.7); sy = Inches(2.3)
    for i, (label, desc) in enumerate(items):
        col, row = i % 2, i // 2
        x = sx + (cw + gx) * col
        y = sy + (ch + gy) * row
        rrect(s, x, y, cw, ch, NAVY_MID)
        text(s, x + Inches(0.35), y + Inches(0.3), cw - Inches(0.7), Inches(0.45),
             label, size=16, bold=True, color=AMBER)
        text(s, x + Inches(0.35), y + Inches(0.85), cw - Inches(0.7), Inches(0.5),
             desc, size=12, color=GRAY_200)

    rrect(s, Inches(0.7), Inches(5.65), Inches(12.15), Inches(1.0), AMBER)
    text(s, Inches(0.7), Inches(5.8), Inches(12.15), Inches(0.4),
         "카톡·엑셀 다 버리지 마세요.",
         size=15, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(6.2), Inches(12.15), Inches(0.4),
         "수플렉스는 그 사이에 흩어진 정보만 모아줍니다.",
         size=15, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
    footer(s, 11, dark=True)


def s_12_onboarding(prs):
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "06. ONBOARDING", size=11, bold=True, color=GRAY_400)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "도입 첫날 30분 — 변화는 가볍게",
         size=30, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.85), Inches(0.6), Inches(0.04), AMBER)

    steps = [
        ("01", "초대 링크 클릭",
         "이름·비밀번호 설정 → 자동으로 회사 가입"),
        ("02", "한국 표준 5묶음 자동 복사",
         "공정·키워드·D-N 룰·체크리스트·견적 가이드가 첫날부터 깔림"),
        ("03", "회사 정보 입력",
         "회사명·로고·연락처. 견적서 PDF 헤더에 자동 반영"),
        ("04", "팀원 초대 (선택)",
         "디자이너·현장팀 역할 부여. 멤버 무제한"),
        ("05", "첫 프로젝트 등록",
         "다음 신규 현장부터 시작. 옛날 프로젝트 다 옮길 필요 X"),
    ]
    y = Inches(2.3)
    for num, title, desc in steps:
        rrect(s, Inches(0.7), y, Inches(12.15), Inches(0.78), GRAY_50)
        circle_num(s, Inches(0.95), y + Inches(0.14), Inches(0.5), num, NAVY_DARK, WHITE)
        text(s, Inches(1.65), y + Inches(0.13), Inches(4.8), Inches(0.45),
             title, size=14, bold=True, color=NAVY_DARK, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(6.5), y + Inches(0.13), Inches(6.3), Inches(0.45),
             desc, size=11, color=GRAY_600, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.88)

    text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         "기존 카톡·엑셀은 그대로 쓰세요. 수플렉스는 그 사이에 흩어진 정보만 모아줍니다.",
         size=12, color=GRAY_600, align=PP_ALIGN.CENTER)
    footer(s, 12)


def s_13_pricing(prs):
    s = add_blank(prs); fill_bg(prs, s, WHITE)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "07. PRICING", size=11, bold=True, color=GRAY_400)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.7),
         "가격 패키지", size=30, bold=True, color=NAVY_DARK)
    rect(s, Inches(0.7), Inches(1.85), Inches(0.6), Inches(0.04), AMBER)

    plans = [
        ("STARTER", "스타터", "₩59,000", "/월",
         "소규모 회사 시작용", "프로젝트·용량 등급별 정원"),
        ("PRO", "프로", "₩99,000", "/월 (베타 ₩79,200)",
         "메인 추천 등급", "프로젝트 무제한 · 멤버 무제한"),
        ("ENT", "엔터프라이즈", "₩179,000", "/월",
         "회사 표준 자유 편집", "+ 프리셋 자체 편집"),
    ]
    cw = Inches(4.0); gx = Inches(0.075); ch = Inches(4.1)
    sx = Inches(0.7); sy = Inches(2.3)
    for i, (key, name, price, unit, desc, feat) in enumerate(plans):
        x = sx + (cw + gx) * i
        hl = (i == 1)
        bg = NAVY_DARK if hl else GRAY_50
        c_key = AMBER if hl else GRAY_400
        c_name = WHITE if hl else NAVY_DARK
        c_price = AMBER if hl else NAVY_DARK
        c_unit = GRAY_200 if hl else GRAY_600
        c_desc = GRAY_200 if hl else GRAY_600
        c_feat = WHITE if hl else GRAY_800

        rrect(s, x, sy, cw, ch, bg)
        text(s, x + Inches(0.35), sy + Inches(0.4), cw - Inches(0.7), Inches(0.4),
             key, size=11, bold=True, color=c_key)
        text(s, x + Inches(0.35), sy + Inches(0.85), cw - Inches(0.7), Inches(0.5),
             name, size=20, bold=True, color=c_name)
        text(s, x + Inches(0.35), sy + Inches(1.65), cw - Inches(0.7), Inches(0.6),
             price, size=26, bold=True, color=c_price)
        text(s, x + Inches(0.35), sy + Inches(2.3), cw - Inches(0.7), Inches(0.4),
             unit, size=11, color=c_unit)
        text(s, x + Inches(0.35), sy + Inches(2.95), cw - Inches(0.7), Inches(0.5),
             desc, size=11, color=c_desc)
        text(s, x + Inches(0.35), sy + Inches(3.5), cw - Inches(0.7), Inches(0.5),
             feat, size=11, bold=True, color=c_feat)
        if hl:
            rrect(s, x + cw - Inches(1.2), sy + Inches(0.35), Inches(0.95), Inches(0.32),
                  AMBER)
            text(s, x + cw - Inches(1.2), sy + Inches(0.38), Inches(0.95), Inches(0.3),
                 "추천", size=10, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)

    text(s, Inches(0.7), Inches(6.65), Inches(12.15), Inches(0.4),
         "베타 회원 평생 20% 할인 — 어느 등급이든 영구 적용",
         size=13, bold=True, color=AMBER_DARK, align=PP_ALIGN.CENTER)
    footer(s, 13)


def s_14_cta(prs):
    s = add_blank(prs); fill_bg(prs, s, NAVY_DARK)
    text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
         "08. NEXT STEP", size=11, bold=True, color=AMBER)
    text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.9),
         "다음 행동", size=36, bold=True, color=WHITE)
    rect(s, Inches(0.7), Inches(1.95), Inches(0.6), Inches(0.04), AMBER)

    actions = [
        ("이메일",
         "데모 · 베타 시범 도입 · 일반 문의",
         "hello@suplex.kr"),
        ("인스타그램",
         "최신 기능 업데이트 · 시범 사례 공유",
         "@suplex.kr"),
        ("카카오톡 채널",
         "법인 등록 후 개설 예정",
         "준비 중"),
    ]
    y = Inches(2.6)
    for label, desc, value in actions:
        rrect(s, Inches(0.7), y, Inches(12.15), Inches(1.05), NAVY_MID)
        text(s, Inches(1.0), y + Inches(0.18), Inches(4), Inches(0.4),
             label, size=18, bold=True, color=AMBER)
        text(s, Inches(1.0), y + Inches(0.6), Inches(8), Inches(0.4),
             desc, size=12, color=GRAY_200)
        text(s, Inches(8.4), y + Inches(0.32), Inches(4.4), Inches(0.5),
             value, size=14, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)
        y += Inches(1.2)

    text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
         "한 사람의 능력이 아니라, 회사 전체의 능력이 올라갑니다.",
         size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    footer(s, 14, dark=True)


# ── 메인 ──
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s_1_title(prs)
    s_2_pain(prs)
    s_3_old_vs_new(prs)
    s_4_solution_overview(prs)
    s_5_solution_1(prs)
    s_6_solution_2(prs)
    s_7_solution_3(prs)
    s_8_solution_4(prs)
    s_9_company_asset(prs)
    s_10_diff(prs)
    s_11_beta(prs)
    s_12_onboarding(prs)
    s_13_pricing(prs)
    s_14_cta(prs)

    out = r"C:\Users\1988k\projects\suplex\docs\sales\소개서.pptx"
    prs.save(out)
    print(f"OK: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
